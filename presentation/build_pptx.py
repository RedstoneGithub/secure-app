from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

# ── Palette ────────────────────────────────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x27, 0x61)
ICE     = RGBColor(0xCA, 0xDC, 0xFC)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT  = RGBColor(0x4A, 0x90, 0xD9)
DARK_TEXT = RGBColor(0x1E, 0x27, 0x61)
LIGHT_BG  = RGBColor(0xF4, 0xF7, 0xFF)
GRAY_TEXT = RGBColor(0x55, 0x65, 0x80)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
RED     = RGBColor(0xDC, 0x26, 0x26)
AMBER   = RGBColor(0xD9, 0x77, 0x06)

BLANK = prs.slide_layouts[6]  # completely blank layout

# ── Helpers ────────────────────────────────────────────────────────────────────
def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h, size, color, bold=False, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = "Calibri"
    return tb

def slide_header(slide, title, light=True):
    """Thin navy top bar with white title text."""
    rect(slide, 0, 0, 13.33, 0.85, NAVY)
    txt(slide, title, 0.4, 0.1, 12.5, 0.65, 26, WHITE, bold=True)

def card(slide, x, y, w, h, title, body, title_color=None, accent_color=None):
    """White card with left accent bar."""
    tc = title_color or NAVY
    ac = accent_color or ACCENT
    rect(slide, x, y, w, h, WHITE)
    # subtle border via a slightly larger gray rect behind
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = ICE
    shape.line.width = Pt(1)
    # left accent bar
    rect(slide, x, y, 0.06, h, ac)
    txt(slide, title, x + 0.15, y + 0.1, w - 0.2, 0.35, 13, tc, bold=True)
    txt(slide, body,  x + 0.15, y + 0.45, w - 0.2, h - 0.55, 11, GRAY_TEXT, wrap=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1, NAVY)

# centre band
rect(s1, 0, 2.5, 13.33, 2.8, RGBColor(0x16, 0x1E, 0x52))

# accent strip
rect(s1, 0, 2.5, 13.33, 0.06, ACCENT)
rect(s1, 0, 5.24, 13.33, 0.06, ACCENT)

txt(s1, "Secure Document Sharing System",
    0.5, 2.75, 12.3, 1.1, 40, WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s1, "CS 419  —  Spring 2026",
    0.5, 3.85, 12.3, 0.6, 20, ICE, align=PP_ALIGN.CENTER)

txt(s1, "Security Engineering Course Project",
    0.5, 1.3, 12.3, 0.6, 16, ICE, align=PP_ALIGN.CENTER, italic=True)

txt(s1, "Flask  •  Fernet Encryption  •  RBAC  •  TLS  •  Audit Logging",
    0.5, 6.2, 12.3, 0.6, 13, RGBColor(0x8A, 0xA8, 0xE8), align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Project Overview
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2, LIGHT_BG)
slide_header(s2, "Project Overview")

# left column — description
txt(s2, "What we built",
    0.4, 1.1, 6.0, 0.5, 16, NAVY, bold=True)
txt(s2,
    "A web application that lets users securely upload, store, share, "
    "and manage documents — built with security as a first principle, "
    "not an afterthought.",
    0.4, 1.6, 5.8, 1.0, 12, GRAY_TEXT)

txt(s2, "Core capabilities",
    0.4, 2.75, 6.0, 0.4, 14, NAVY, bold=True)

features = [
    "Role-based access control  (Admin / User / Guest)",
    "End-to-end encryption at rest with Fernet (AES-128-CBC)",
    "Document versioning and restore",
    "Secure sharing with Editor / Viewer permissions",
    "Full audit trail across two structured log files",
]
y = 3.2
for f in features:
    rect(s2, 0.4, y + 0.12, 0.06, 0.18, ACCENT)
    txt(s2, f, 0.56, y, 5.7, 0.38, 11.5, DARK_TEXT)
    y += 0.38

# right column — stat callouts
stats = [
    ("3",    "User Roles"),
    ("AES-128", "Encryption"),
    ("7",    "Security Headers"),
    ("130+", "Automated Tests"),
]
sx, sy = 7.2, 1.1
for i, (num, label) in enumerate(stats):
    cx = sx + (i % 2) * 2.9
    cy = sy + (i // 2) * 1.6
    rect(s2, cx, cy, 2.6, 1.35, WHITE)
    # border
    b = s2.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(2.6), Inches(1.35))
    b.fill.background(); b.line.color.rgb = ICE; b.line.width = Pt(1)
    txt(s2, num,   cx + 0.1, cy + 0.1,  2.4, 0.7,  32, ACCENT,    bold=True, align=PP_ALIGN.CENTER)
    txt(s2, label, cx + 0.1, cy + 0.75, 2.4, 0.45, 11, GRAY_TEXT, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Security Architecture
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3, LIGHT_BG)
slide_header(s3, "Security Architecture")

# architecture flow: Browser → TLS → Flask → Storage
layers = [
    ("Browser",        "HTML5 / JS\nUser Interface"),
    ("TLS / HTTPS",    "Self-signed cert\nPort 5000"),
    ("Flask App",      "Routes · Auth\nRBAC · Validation"),
    ("Encrypted Store","Fernet .enc files\nJSON sessions"),
]
box_w, box_h = 2.2, 1.3
total_w = len(layers) * box_w + (len(layers) - 1) * 0.55
start_x = (13.33 - total_w) / 2

for i, (title, sub) in enumerate(layers):
    bx = start_x + i * (box_w + 0.55)
    by = 1.5
    color = NAVY if i in (0, 3) else ACCENT
    rect(s3, bx, by, box_w, box_h, color)
    txt(s3, title, bx, by + 0.1, box_w, 0.5, 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    txt(s3, sub,   bx, by + 0.6, box_w, 0.6, 10, ICE,   align=PP_ALIGN.CENTER)
    # arrow
    if i < len(layers) - 1:
        ax = bx + box_w + 0.05
        ay = by + box_h / 2 - 0.04
        rect(s3, ax, ay, 0.45, 0.08, GRAY_TEXT)

# security layers below
txt(s3, "Security Controls Applied at Each Layer",
    0.5, 3.15, 12.3, 0.45, 14, NAVY, bold=True, align=PP_ALIGN.CENTER)

controls = [
    ("Authentication",  "bcrypt · lockout · rate limiting"),
    ("Authorization",   "RBAC decorators · ownership checks"),
    ("Input Handling",  "Whitelist validation · path traversal prevention"),
    ("Data Protection", "Fernet encryption · secure deletion"),
    ("Transport",       "HTTPS redirect · HSTS · SameSite cookies"),
    ("Observability",   "security.log · access.log · structured JSON"),
]
cw = 3.8
cx_start = 0.35
for i, (name, detail) in enumerate(controls):
    col = i % 3
    row = i // 3
    cx = cx_start + col * (cw + 0.3)
    cy = 3.7 + row * 1.4
    rect(s3, cx, cy, cw, 1.2, WHITE)
    b2 = s3.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(cw), Inches(1.2))
    b2.fill.background(); b2.line.color.rgb = ICE; b2.line.width = Pt(1)
    rect(s3, cx, cy, cw, 0.06, ACCENT)
    txt(s3, name,   cx + 0.15, cy + 0.12, cw - 0.2, 0.35, 12, NAVY, bold=True)
    txt(s3, detail, cx + 0.15, cy + 0.48, cw - 0.2, 0.6,  10, GRAY_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Key Security Implementations
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4, LIGHT_BG)
slide_header(s4, "Key Security Implementations")

impls = [
    ("Authentication",
     "bcrypt (cost 12) · 5-attempt lockout · IP rate limiting · password policy enforcement"),
    ("Encryption at Rest",
     "Fernet (AES-128-CBC + HMAC-SHA256) · all documents encrypted before disk write · key auto-generated"),
    ("Role-Based Access Control",
     "3 roles: Admin / User / Guest · decorator-enforced · editor vs viewer distinction on shared docs"),
    ("Session Management",
     "cryptographic tokens (secrets.token_urlsafe) · 30 min timeout · HttpOnly + Secure + SameSite=Strict"),
    ("Input Validation",
     "UUID regex on doc IDs · extension + MIME whitelist · 10 MB limit · Jinja2 auto-escaping for XSS"),
    ("Audit Logging",
     "security.log for auth/security events · access.log for data access · structured JSON with IP + timestamp"),
]

iw, ih = 5.9, 1.45
for i, (name, detail) in enumerate(impls):
    col = i % 2
    row = i // 2
    ix = 0.35 + col * (iw + 0.7)
    iy = 1.05 + row * (ih + 0.18)
    rect(s4, ix, iy, iw, ih, WHITE)
    b = s4.shapes.add_shape(1, Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    b.fill.background(); b.line.color.rgb = ICE; b.line.width = Pt(1)
    rect(s4, ix, iy, 0.07, ih, ACCENT)
    txt(s4, name,   ix + 0.2, iy + 0.1,  iw - 0.25, 0.38, 13, NAVY, bold=True)
    txt(s4, detail, ix + 0.2, iy + 0.52, iw - 0.25, ih - 0.6, 10.5, GRAY_TEXT)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Penetration Testing Results
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
bg(s5, LIGHT_BG)
slide_header(s5, "Penetration Testing Results")

# severity summary boxes
severities = [
    ("Critical", "0", RED),
    ("High",     "0", RGBColor(0xEA, 0x58, 0x0C)),
    ("Medium",   "0", AMBER),
    ("Low",      "2", RGBColor(0x16, 0x78, 0xD4)),
    ("Info",     "1", GRAY_TEXT),
]
bw = 2.0
bx_start = (13.33 - len(severities) * bw - (len(severities)-1)*0.3) / 2
for i, (label, count, color) in enumerate(severities):
    bx = bx_start + i * (bw + 0.3)
    by = 1.05
    rect(s5, bx, by, bw, 1.35, WHITE)
    b = s5.shapes.add_shape(1, Inches(bx), Inches(by), Inches(bw), Inches(1.35))
    b.fill.background(); b.line.color.rgb = color; b.line.width = Pt(2)
    rect(s5, bx, by, bw, 0.07, color)
    txt(s5, count, bx, by + 0.1,  bw, 0.7,  36, color, bold=True, align=PP_ALIGN.CENTER)
    txt(s5, label, bx, by + 0.85, bw, 0.38, 12, GRAY_TEXT, align=PP_ALIGN.CENTER)

# findings table
txt(s5, "Notable Findings", 0.5, 2.65, 8.0, 0.4, 14, NAVY, bold=True)

findings = [
    ("LOW",  "Username enumeration",        "Registration reveals whether username or email already exists separately."),
    ("LOW",  "No rate limit on /change-password", "Password change endpoint lacks the throttling applied to /login."),
    ("INFO", "Secure deletion on SSDs",     "Byte-overwrite before delete is best-effort only on modern storage."),
]

headers = ["Severity", "Finding", "Description"]
header_widths = [1.2, 3.2, 7.8]
row_y = 3.15
hx = 0.35
# header row
rect(s5, 0.35, row_y, 12.6, 0.42, NAVY)
col_x = hx
for hw, htext in zip(header_widths, headers):
    txt(s5, htext, col_x + 0.1, row_y + 0.08, hw - 0.1, 0.3, 11, WHITE, bold=True)
    col_x += hw

row_y += 0.42
for sev, name, desc in findings:
    color = RGBColor(0x16, 0x78, 0xD4) if sev == "LOW" else GRAY_TEXT
    rect(s5, 0.35, row_y, 12.6, 0.52, WHITE)
    b = s5.shapes.add_shape(1, Inches(0.35), Inches(row_y), Inches(12.6), Inches(0.52))
    b.fill.background(); b.line.color.rgb = ICE; b.line.width = Pt(0.5)
    col_x = hx
    txt(s5, sev,  col_x + 0.1, row_y + 0.1, 1.1,  0.35, 10, color, bold=True); col_x += 1.2
    txt(s5, name, col_x + 0.1, row_y + 0.1, 3.1,  0.35, 10, DARK_TEXT, bold=True); col_x += 3.2
    txt(s5, desc, col_x + 0.1, row_y + 0.05, 7.6, 0.42, 9.5, GRAY_TEXT)
    row_y += 0.54

txt(s5, "All critical, high, and medium attack vectors were mitigated successfully.",
    0.35, row_y + 0.2, 12.6, 0.4, 12, GREEN, bold=True, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Challenges & Lessons Learned
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
bg(s6, LIGHT_BG)
slide_header(s6, "Challenges & Lessons Learned")

# two columns
col_titles = ["Challenges", "Lessons Learned"]
col_items = [
    [
        ("File-based storage security",
         "Ensuring encrypted JSON files on disk couldn't be tampered with or accessed outside the app required careful path validation and Fernet HMAC verification."),
        ("Balancing usability and security",
         "Strict SameSite=Strict cookies and HTTPS-only enforcement created friction during local development, requiring debug-mode exceptions."),
        ("Secure deletion limitations",
         "Discovered that byte-overwriting files does not guarantee erasure on SSDs — shifted reliance to encryption-key management as the primary deletion defense."),
    ],
    [
        ("Defence in depth matters",
         "No single control is enough. Combining RBAC, encryption, validation, and logging means an attacker must defeat multiple independent layers."),
        ("Test the negative cases",
         "Security testing is mostly about proving things don't work when they shouldn't — access denied, lockout enforced, bad input rejected."),
        ("Secrets in version control are a real risk",
         "Keeping encryption keys and TLS certs in git — even for convenience — is a practice that must be avoided in any non-demo deployment."),
    ],
]

cw = 5.9
for col, (ctitle, items) in enumerate(zip(col_titles, col_items)):
    cx = 0.35 + col * (cw + 0.75)
    rect(s6, cx, 1.0, cw, 0.5, NAVY)
    txt(s6, ctitle, cx + 0.2, 1.05, cw - 0.25, 0.38, 14, WHITE, bold=True)
    iy = 1.6
    for title, detail in items:
        rect(s6, cx, iy, cw, 1.55, WHITE)
        b = s6.shapes.add_shape(1, Inches(cx), Inches(iy), Inches(cw), Inches(1.55))
        b.fill.background(); b.line.color.rgb = ICE; b.line.width = Pt(1)
        rect(s6, cx, iy, 0.07, 1.55, ACCENT)
        txt(s6, title,  cx + 0.2, iy + 0.1,  cw - 0.25, 0.38, 12, NAVY, bold=True)
        txt(s6, detail, cx + 0.2, iy + 0.5,  cw - 0.25, 1.0,  10, GRAY_TEXT)
        iy += 1.65

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Conclusion
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
bg(s7, NAVY)
rect(s7, 0, 2.4, 13.33, 0.06, ACCENT)
rect(s7, 0, 5.0, 13.33, 0.06, ACCENT)

txt(s7, "Conclusion", 0.5, 1.3, 12.3, 0.7, 28, ICE, bold=True, align=PP_ALIGN.CENTER)

summary = (
    "We built a secure document sharing system that demonstrates production-grade security "
    "practices across every layer — from bcrypt password hashing and Fernet encryption at rest, "
    "to strict RBAC, input validation, session security, and structured audit logging.\n\n"
    "Penetration testing confirmed no critical, high, or medium vulnerabilities. "
    "The two low-severity findings identified are known and have clear remediation paths."
)
txt(s7, summary, 1.0, 2.6, 11.3, 2.2, 14, WHITE, align=PP_ALIGN.CENTER)

takeaways = [
    "Security is built in layers — no single control is sufficient",
    "Test what shouldn't work as rigorously as what should",
    "Encryption + key management is stronger than secure deletion alone",
]
ty = 5.2
for t in takeaways:
    rect(s7, 2.8, ty + 0.12, 0.08, 0.18, ACCENT)
    txt(s7, t, 3.0, ty, 7.5, 0.38, 12, ICE)
    ty += 0.38

txt(s7, "CS 419  —  Spring 2026", 0.5, 7.0, 12.3, 0.35, 11,
    RGBColor(0x8A, 0xA8, 0xE8), align=PP_ALIGN.CENTER, italic=True)

# ── Save ───────────────────────────────────────────────────────────────────────
out = r"C:\Users\nilay\OneDrive\Documents\CS419-FinalProj\presentation\CS419_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
